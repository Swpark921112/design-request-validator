"""PPT 기획서 파싱 모듈 (python-pptx)

PPT 기획서 템플릿은 카드 기반 레이아웃을 사용한다:
- 카드 배경 (AUTO_SHAPE 사각형)
- 레이블 텍스트박스 ("요청자 *")
- 값 텍스트박스 ("이름 - 소속팀")  ← 이것이 플레이스홀더/실제 값
- 예시 텍스트박스 ("예) 박상욱 - B2D")

각 카드 내에서 레이블 → 값 → 예시 순서로 위에서 아래로 배치된다.
"""

from __future__ import annotations

from dataclasses import dataclass, field

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from config import PLACEHOLDER_KEYWORDS, REQUIRED_FIELDS, SLIDE_INDEX


@dataclass
class TextShape:
    """텍스트가 있는 Shape 정보"""
    text: str
    left: int
    top: int
    width: int
    height: int


@dataclass
class CardField:
    """카드에서 추출한 필드 정보"""
    label: str       # 레이블 ("요청자 *")
    value: str       # 사용자 입력 값 또는 플레이스홀더
    example: str     # 예시 텍스트


@dataclass
class SlideContent:
    """슬라이드에서 추출한 텍스트 정보"""
    index: int
    title: str = ""
    texts: list[str] = field(default_factory=list)
    fields: dict[str, str] = field(default_factory=dict)
    cards: list[CardField] = field(default_factory=list)


@dataclass
class ParsedPPTX:
    """파싱된 PPT 기획서 전체"""
    slide_count: int = 0
    slides: dict[str, SlideContent] = field(default_factory=dict)
    all_text: str = ""
    raw_slides: list[SlideContent] = field(default_factory=list)


def _get_text_shapes(slide) -> list[TextShape]:
    """슬라이드의 텍스트 shape를 좌표 포함하여 추출한다."""
    shapes = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                shapes.append(TextShape(
                    text=text,
                    left=shape.left,
                    top=shape.top,
                    width=shape.width,
                    height=shape.height,
                ))
    return shapes


def _get_card_bounds(slide) -> list[dict]:
    """카드 배경 사각형(AUTO_SHAPE)의 좌표를 추출한다."""
    cards = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            # 일정 크기 이상이고 텍스트가 없는 사각형 = 카드 배경
            if shape.width > 1000000 and shape.height > 400000:
                if not shape.has_text_frame or not shape.text_frame.text.strip():
                    cards.append({
                        "left": shape.left,
                        "top": shape.top,
                        "right": shape.left + shape.width,
                        "bottom": shape.top + shape.height,
                    })
    return cards


def _shapes_in_card(text_shapes: list[TextShape], card: dict) -> list[TextShape]:
    """카드 영역 안에 있는 텍스트 shape를 top 순서로 반환한다."""
    inside = []
    for ts in text_shapes:
        # shape의 중심이 카드 안에 있는지
        cx = ts.left + ts.width // 2
        cy = ts.top + ts.height // 2
        if card["left"] <= cx <= card["right"] and card["top"] <= cy <= card["bottom"]:
            inside.append(ts)
    inside.sort(key=lambda s: s.top)
    return inside


def _extract_card_fields(slide) -> list[CardField]:
    """슬라이드에서 카드 기반 필드를 추출한다."""
    text_shapes = _get_text_shapes(slide)
    cards = _get_card_bounds(slide)
    result = []

    for card in cards:
        inside = _shapes_in_card(text_shapes, card)
        if len(inside) >= 2:
            label = inside[0].text
            value = inside[1].text
            example = inside[2].text if len(inside) >= 3 else ""
            result.append(CardField(label=label, value=value, example=example))
        elif len(inside) == 1:
            result.append(CardField(label=inside[0].text, value="", example=""))

    return result


def _extract_all_texts(slide) -> list[str]:
    """슬라이드의 모든 shape에서 텍스트를 추출한다."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    texts.append(text)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    text = cell.text.strip()
                    if text:
                        texts.append(text)
    return texts


# 플레이스홀더 감지용 추가 패턴 (값 텍스트박스에 남아있는 안내 문구)
PLACEHOLDER_VALUE_PATTERNS = [
    "이름 - 소속팀",
    "피드백 취합",
    "프로젝트 / 캠페인",
    "yyyy-mm-dd",
    "입력",
    "입력하세요",
    "선택하세요",
    "명시",
    "기재",
    "붙여넣",
    # 선택지 나열 형태 (드롭다운/안내 문구)
    "배너 / 리플렛",
    "ai / pdf / png",
    "가로 x 세로",
    "발주처, 제작 납기",
    "특정 도구 필요 시",
    "인쇄물/패키지인 경우 필수",
]


def is_placeholder(value: str) -> bool:
    """값이 플레이스홀더(미작성) 상태인지 판별한다."""
    if not value:
        return True
    lower = value.lower().strip()
    for kw in PLACEHOLDER_KEYWORDS:
        if kw in lower:
            return True
    for pattern in PLACEHOLDER_VALUE_PATTERNS:
        if pattern in lower:
            return True
    return False


def parse_pptx(file_path: str) -> ParsedPPTX:
    """PPT 기획서를 파싱하여 구조화된 데이터를 반환한다."""
    prs = Presentation(file_path)
    result = ParsedPPTX(slide_count=len(prs.slides))

    all_texts = []

    for i, slide in enumerate(prs.slides):
        texts = _extract_all_texts(slide)
        cards = _extract_card_fields(slide)
        title = ""
        if slide.shapes.title:
            title = slide.shapes.title.text.strip()

        # 카드 기반 필드를 dict로 변환 (label → value)
        fields = {}
        for card in cards:
            fields[card.label] = card.value

        sc = SlideContent(index=i, title=title, texts=texts, fields=fields, cards=cards)
        result.raw_slides.append(sc)
        all_texts.extend(texts)

    result.all_text = "\n".join(all_texts)

    # 슬라이드 이름 매핑
    for name, idx in SLIDE_INDEX.items():
        if idx < len(result.raw_slides):
            result.slides[name] = result.raw_slides[idx]

    return result


def check_required_fields(parsed: ParsedPPTX) -> list[dict]:
    """필수 필드 작성 여부를 검사한다.

    카드 기반 레이아웃에서 레이블로 필드를 찾고,
    해당 카드의 값이 플레이스홀더인지 실제 값인지 판별한다.

    Returns:
        list of {slide, field, status, value, message}
    """
    results = []
    for slide_name, required in REQUIRED_FIELDS.items():
        slide = parsed.slides.get(slide_name)
        if not slide:
            for f in required:
                results.append({
                    "slide": slide_name,
                    "field": f,
                    "status": "fail",
                    "value": "",
                    "message": f"슬라이드 '{slide_name}'를 찾을 수 없습니다.",
                })
            continue

        for req_field in required:
            # 카드에서 매칭
            matched_card = None
            for card in slide.cards:
                if req_field in card.label:
                    matched_card = card
                    break

            # fields dict에서 매칭 (카드에서 못 찾은 경우)
            if matched_card is None:
                for key, val in slide.fields.items():
                    if req_field in key:
                        matched_card = CardField(label=key, value=val, example="")
                        break

            if matched_card is None:
                results.append({
                    "slide": slide_name,
                    "field": req_field,
                    "status": "fail",
                    "value": "",
                    "message": f"'{req_field}' 항목을 찾을 수 없습니다.",
                })
            elif is_placeholder(matched_card.value):
                results.append({
                    "slide": slide_name,
                    "field": req_field,
                    "status": "fail",
                    "value": matched_card.value,
                    "message": f"'{req_field}' 항목이 미작성 상태입니다. 실제 내용을 작성해주세요.",
                })
            else:
                results.append({
                    "slide": slide_name,
                    "field": req_field,
                    "status": "pass",
                    "value": matched_card.value,
                    "message": f"'{req_field}' 작성 완료: {matched_card.value[:50]}",
                })

    return results
